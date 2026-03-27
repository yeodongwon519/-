import json
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import math

# ── 색상 정의 ──────────────────────────────────────────
BG_COLOR      = RGBColor(0x1C, 0x1C, 0x1C)   # 거의 검정
CARD_COLOR    = RGBColor(0x2C, 0x2C, 0x2C)   # 다크 그레이 카드
ORANGE        = RGBColor(0xE8, 0x64, 0x3A)   # 주황 강조
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xAA, 0xAA, 0xAA)

# ── 슬라이드 크기 (9:16 세로형) ───────────────────────
SLIDE_W = Inches(7.5)
SLIDE_H = Inches(13.33)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, text, left, top, width, height,
                 font_size=28, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, word_wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Pretendard"   # 없으면 시스템 기본 폰트 사용
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color, radius=0.15):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE = 5
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()   # 테두리 제거
    try:
        shape.adjustments[0] = radius
    except (IndexError, Exception):
        pass
    return shape


def add_circle(slide, cx, cy, r, fill_color, text, font_size=22):
    left = cx - r
    top  = cy - r
    size = r * 2
    shape = slide.shapes.add_shape(9, left, top, size, size)  # 9 = OVAL
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = WHITE
    return shape


def make_opening_slide(prs, scene):
    """Scene 1 - 오프닝: 제목만 크게"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_COLOR)

    # 주황 강조 라인 (상단)
    line = slide.shapes.add_shape(1, Inches(0), Inches(0.5), SLIDE_W, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    # 중앙 메인 텍스트
    add_text_box(
        slide, scene["overlay"],
        Inches(0.5), Inches(4.5), Inches(6.5), Inches(2),
        font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )

    # 서브텍스트 (TTS 첫 문장 요약)
    short = scene["tts"][:30] + "..."
    add_text_box(
        slide, short,
        Inches(0.5), Inches(6.8), Inches(6.5), Inches(1.5),
        font_size=22, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER
    )

    # 주황 강조 라인 (하단)
    line2 = slide.shapes.add_shape(1, Inches(0), Inches(12.77), SLIDE_W, Inches(0.06))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ORANGE
    line2.line.fill.background()


def make_content_slide(prs, scene, number):
    """Scene 2~4 - 본문: 번호 + 카드 리스트"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_COLOR)

    # 상단 주황 라인
    line = slide.shapes.add_shape(1, Inches(0), Inches(0.5), SLIDE_W, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    # 슬라이드 타입 레이블 (상단 좌측)
    add_text_box(
        slide, scene["type"].upper(),
        Inches(0.5), Inches(0.7), Inches(3), Inches(0.5),
        font_size=16, bold=True, color=ORANGE
    )

    # 제목
    add_text_box(
        slide, scene["overlay"],
        Inches(0.5), Inches(1.4), Inches(6.5), Inches(1.2),
        font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT
    )

    # 구분선
    sep = slide.shapes.add_shape(1, Inches(0.5), Inches(2.9), Inches(6.5), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = RGBColor(0x44, 0x44, 0x44)
    sep.line.fill.background()

    # TTS 텍스트를 문장으로 분리해서 카드로 표시
    sentences = [s.strip() for s in scene["tts"].replace(".", ".|").split("|") if s.strip()]

    card_top = Inches(3.2)
    card_h   = Inches(1.5)
    gap      = Inches(0.25)
    card_l   = Inches(1.2)
    card_w   = Inches(5.8)
    circle_r = Inches(0.35)

    for i, sent in enumerate(sentences[:3]):   # 최대 3개 카드
        top = card_top + i * (card_h + gap)

        # 카드 배경
        add_rounded_rect(slide, card_l, top, card_w, card_h, CARD_COLOR)

        # 번호 원
        add_circle(
            slide,
            Inches(0.6), top + card_h / 2,
            circle_r, ORANGE, str(i + 1), font_size=20
        )

        # 카드 텍스트
        add_text_box(
            slide, sent,
            card_l + Inches(0.2), top + Inches(0.15),
            card_w - Inches(0.4), card_h - Inches(0.3),
            font_size=18, bold=False, color=WHITE, word_wrap=True
        )

    # 하단 주황 라인
    line2 = slide.shapes.add_shape(1, Inches(0), Inches(12.77), SLIDE_W, Inches(0.06))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ORANGE
    line2.line.fill.background()


def make_closing_slide(prs, scene):
    """Scene 5 - 클로징 CTA"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG_COLOR)

    # 주황 라인 상단
    line = slide.shapes.add_shape(1, Inches(0), Inches(0.5), SLIDE_W, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    # "SUMMARY" 레이블
    add_text_box(
        slide, "SUMMARY",
        Inches(0), Inches(2.5), SLIDE_W, Inches(0.8),
        font_size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER
    )

    # 메인 CTA 텍스트
    add_text_box(
        slide, scene["overlay"],
        Inches(0.5), Inches(3.5), Inches(6.5), Inches(2),
        font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER
    )

    # 주황 구분선
    sep = slide.shapes.add_shape(1, Inches(2), Inches(6.2), Inches(3.5), Inches(0.06))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ORANGE
    sep.line.fill.background()

    # 본문 요약
    add_text_box(
        slide, scene["tts"],
        Inches(0.5), Inches(6.5), Inches(6.5), Inches(3),
        font_size=19, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.CENTER, word_wrap=True
    )

    # 하단 주황 라인
    line2 = slide.shapes.add_shape(1, Inches(0), Inches(12.77), SLIDE_W, Inches(0.06))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ORANGE
    line2.line.fill.background()


# ── 메인 실행 ──────────────────────────────────────────
def create_ppt(script_path, output_path):
    with open(script_path, "r", encoding="utf-8") as f:
        script = json.load(f)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, scene in enumerate(script["scenes"]):
        if i == 0:
            make_opening_slide(prs, scene)
        elif i == len(script["scenes"]) - 1:
            make_closing_slide(prs, scene)
        else:
            make_content_slide(prs, scene, i)

    prs.save(output_path)
    print(f"✅ PPT 저장 완료: {output_path}")


if __name__ == "__main__":
    base = os.path.dirname(os.path.abspath(__file__))
    script_path = os.path.join(base, "output/01_script/rsi_script.json")
    output_path = os.path.join(base, "output/03_ppt/rsi_slides.pptx")
    create_ppt(script_path, output_path)
